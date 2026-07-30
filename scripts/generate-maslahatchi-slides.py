# -*- coding: utf-8 -*-
"""7 betli taqdimot: Maktab maslahatchisi + Kelajak Markazi ilovasi."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent.parent / "docs" / "Maktab-maslahatchi-Kelajak-Markazi.pptx"

PURPLE = RGBColor(0x5B, 0x21, 0xB6)
BLUE = RGBColor(0x1E, 0x40, 0xAF)
DARK = RGBColor(0x1E, 0x29, 0x3B)
GRAY = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x7C, 0x3A, 0xED)
GREEN = RGBColor(0x05, 0x96, 0x69)


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header_bar(slide, prs):
    shape = slide.shapes.add_shape(
        1, 0, 0, prs.slide_width, Inches(0.12)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PURPLE
    shape.line.fill.background()


def add_footer(slide, prs, text, num):
    box = slide.shapes.add_textbox(
        Inches(0.5), prs.slide_height - Inches(0.45), Inches(9), Inches(0.35)
    )
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    num_box = slide.shapes.add_textbox(
        prs.slide_width - Inches(1.2), prs.slide_height - Inches(0.45), Inches(0.7), Inches(0.35)
    )
    np = num_box.text_frame.paragraphs[0]
    np.text = num
    np.font.size = Pt(10)
    np.font.color.rgb = GRAY
    np.alignment = PP_ALIGN.RIGHT


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(9), Inches(1.1))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BLUE
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(13)
        p2.font.color.rgb = GRAY
        p2.space_before = Pt(8)


def add_bullets(slide, items, left=0.55, top=1.55, width=9.0, size=15, color=DARK):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.2))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            title, body = item
            p.text = title
            p.font.size = Pt(size)
            p.font.bold = True
            p.font.color.rgb = ACCENT
            p.space_before = Pt(10 if i else 0)
            p2 = tf.add_paragraph()
            p2.text = f"   {body}"
            p2.font.size = Pt(size - 1)
            p2.font.color.rgb = color
            p2.level = 0
        else:
            p.text = item
            p.font.size = Pt(size)
            p.font.color.rgb = color
            p.space_before = Pt(6 if i else 0)
            p.level = 0


def add_two_columns(slide, left_items, right_items):
    add_bullets(slide, left_items, left=0.55, top=1.55, width=4.3, size=14)
    add_bullets(slide, right_items, left=5.15, top=1.55, width=4.3, size=14)


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    footer = "Kelajak Markazi · Maktab maslahatchisi · 2026"

    # --- Slayd 1 ---
    s1 = prs.slides.add_slide(blank)
    set_bg(s1, WHITE)
    add_header_bar(s1, prs)
    add_title(
        s1,
        "Maktab maslahatchisining asosiy vazifalari",
        "O‘zbekiston Respublikasi maktab maslahatchisi — o‘quvchilarning ijtimoiy, kasbiy va madaniy rivojlanishini muvofiqlashtiruvchi mutaxassis",
    )
    add_two_columns(
        s1,
        [
            ("1. Prezident iqtidorli farzandlari", "Milliy dastur mazmunini yoritish va targ‘ib qilish · PF-86 farmon (2025-yil 15-may)"),
            ("2. Kelajak markazi tarmoq to‘garaklari", "Madaniyat, san’at, kitobxonlik, sport, robototexnika · 343-son PQ (2019-yil 22-aprel)"),
        ],
        [
            ("3. Tarix va qadriyatlar", "Tarixiy obidalar, madaniy yodgorliklar, milliy qadriyatlar · Samarqand, Buxoro, Xiva; virtual turlar"),
            ("4. Yetakchilik va kengash", "Loyihalar, klublar, debatlar, «O‘quvchilar kengashi» · 6–11 sinflar, saylov jarayoni"),
        ],
    )
    note = s1.shapes.add_textbox(Inches(0.55), Inches(6.35), Inches(9), Inches(0.5))
    note.text_frame.paragraphs[0].text = "Maslahatchi barcha sinf o‘quvchilari bilan ishlaydi — iqtidorli bolalar dasturidan kasb tanlovigacha."
    note.text_frame.paragraphs[0].font.size = Pt(12)
    note.text_frame.paragraphs[0].font.color.rgb = BLUE
    add_footer(s1, prs, footer, "1 / 7")

    # --- Slayd 2 ---
    s2 = prs.slides.add_slide(blank)
    set_bg(s2, WHITE)
    add_header_bar(s2, prs)
    add_title(s2, "Asosiy funksiyalar (davomi)", "Loyihalar, ma’lumotlar bazasi va kasbga yo‘naltirish")
    add_bullets(
        s2,
        [
            ("Madaniyat · san’at · sport · ekologiya", "Estetik did, jismoniy faollik, kitobxonlik va tabiat bilimini oshirish loyihalari"),
            ("Ijtimoiy portfolio", "Kelajak markazlari, musiqa/san’at maktablari, sport maktablari, Yosh markazlari, kutubxonalar, muzeylar — ixtiyoriy ishtirok"),
            ("Kasbga yo‘naltirish", "So‘rovnomalar; 7–9 sinflar: «Mening kelajakdagi kasbim» · psixolog, kasb vakillari, ekskursiyalar"),
            "Portfolio kuzatiladi: tarmoq to‘garaklari · xalqaro musobaqalar · fan olimpiadalari · IELTS/SAT · kengash a’zoligi",
            "Klublar: Turon teatr · Iqtidor ansambli · Jadidlar izidan · Eco-schools · Xorijiy tillar · Debat · Raqamli avlod qizlari · Inklyuziv klublar",
        ],
        top=1.45,
        size=14,
    )
    add_footer(s2, prs, footer, "2 / 7")

    # --- Slayd 3 ---
    s3 = prs.slides.add_slide(blank)
    set_bg(s3, WHITE)
    add_header_bar(s3, prs)
    add_title(s3, "Kasbga yo‘naltirish tamoyillari", "1-sinfdan 11-sinfgacha bosqichma-bosqich tayyorgarlik")
    add_two_columns(
        s3,
        [
            "Yoshga mos: 1–4 o‘yin · 5–7 loyiha/test · 8–11 amaliy tajriba",
            "Amaliy tajriba: ustaxonalar, laboratoriyalar, mahorat darslari",
            "Qiziqish va qobiliyat: testlar, so‘rovnomalar, kuzatuv",
        ],
        [
            "Kasb vakillari: ishlab chiqarish tashriflari, «bir kun kasbda»",
            "Mehnat bozori tahlili: mintaqa ehtiyojlari, talab yuqori kasblar",
            "Ongli tanlov: 11-sinf tajriba va o‘z-o‘zini tahlil asosida",
        ],
    )
    extra = s3.shapes.add_textbox(Inches(0.55), Inches(5.5), Inches(9), Inches(1.2))
    tf = extra.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Prezident iqtidorli farzandlari: PM hamkorligi · TOP-300 universitetlar · MOCK imtihonlar"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p2 = tf.add_paragraph()
    p2.text = "Tarmoq to‘garaklari: Madaniyat · Robototexnika · San’at · Kitobxonlik · Sport · 9-sinfning 50% ini texnikumga yo‘naltirish"
    p2.font.size = Pt(13)
    p2.font.color.rgb = DARK
    p2.space_before = Pt(8)
    add_footer(s3, prs, footer, "3 / 7")

    # --- Slayd 4 ---
    s4 = prs.slides.add_slide(blank)
    set_bg(s4, RGBColor(0xF8, 0xFA, 0xFC))
    add_header_bar(s4, prs)
    add_title(
        s4,
        "Kelajak Markazi ilovasi",
        "Web + mobil platforma — markaz, maktab va ota-onani bitta raqamli tizimda birlashtiradi",
    )
    stats = s4.shapes.add_textbox(Inches(0.55), Inches(1.5), Inches(9), Inches(0.8))
    sp = stats.text_frame.paragraphs[0]
    sp.text = "Web boshqaruv  ·  Mobil APK  ·  Markaz API  ·  1–11 sinf qamrovi"
    sp.font.size = Pt(18)
    sp.font.bold = True
    sp.font.color.rgb = GREEN
    add_two_columns(
        s4,
        [
            "Nima beradi?",
            "O‘quvchi, to‘garak, to‘lov, davomat — bir joyda",
            "Ota-ona kabineti: to‘lov, xabarlar, faoliyat",
            "Onlayn ro‘yxatdan o‘tish va tuman ko‘rinishi",
        ],
        [
            "Nega maslahatchi kerak?",
            "8+ rasmiy vazifa ilovada bajariladi",
            "Ijtimoiy portfolio elektron shaklda",
            "Rahbariyat real vaqtda natijani ko‘radi",
        ],
    )
    add_footer(s4, prs, footer, "4 / 7")

    # --- Slayd 5 ---
    s5 = prs.slides.add_slide(blank)
    set_bg(s5, RGBColor(0xF8, 0xFA, 0xFC))
    add_header_bar(s5, prs)
    add_title(s5, "Maslahatchi profili", "Har bir maktab maslahatchisi uchun alohida shaxsiy kabinet")
    add_bullets(
        s5,
        [
            ("Akkaunt", "Maktab, F.I.Sh., telefon, mas’ul sinflar — bitta profil"),
            ("Topshiriqlar paneli", "PF-86, 343-PQ va markaz vazifalari — muddat, holat, tasdiq"),
            ("Ijtimoiy portfolio", "Har bir o‘quvchining ishtirokini elektron kiritish va yangilash"),
            "Natija: shaxsiy reyting, statistika, maktab · markaz · klub ma’lumoti bir joyda",
            "Maslahatchi barcha yosh guruhlarini bitta ilovadan boshqaradi — daftar va Excel kerak emas",
        ],
        top=1.45,
        size=15,
    )
    add_footer(s5, prs, footer, "5 / 7")

    # --- Slayd 6 ---
    s6 = prs.slides.add_slide(blank)
    set_bg(s6, RGBColor(0xF8, 0xFA, 0xFC))
    add_header_bar(s6, prs)
    add_title(s6, "Reyting va topshiriqlar tizimi", "Har bir maslahatchi reytingli tartibda ishlaydi")
    add_bullets(
        s6,
        [
            "Reyting mezonlari: Topshiriq bajarilishi 40% · Portfolio 25% · O‘quvchi qamrovi 20% · Loyiha/tadbir 15%",
            "1. Berish — Markaz/tuman standart topshiriq yuboradi",
            "2. Bajarish — Maslahatchi foto, ro‘yxat, hisobot yuklaydi",
            "3. Tasdiqlash — Bosh mutaxassis ball qo‘yadi",
            "4. Reyting — Umumiy ball yangilanadi, tuman/viloyat jadvalida o‘rin",
            "Topshiriqlar: to‘garak yo‘naltirish · kengash saylovi · kasb so‘rovnomasi · ekskursiya · portfolio yangilanishi",
        ],
        top=1.45,
        size=14,
    )
    add_footer(s6, prs, footer, "6 / 7")

    # --- Slayd 7 ---
    s7 = prs.slides.add_slide(blank)
    set_bg(s7, WHITE)
    add_header_bar(s7, prs)
    add_title(s7, "Barcha yosh toifasini qamrab olish", "1-sinfdan 11-sinfgacha — har bosqich uchun alohida modul")
    add_bullets(
        s7,
        [
            "1–4 sinf: o‘yin orqali kasb tanishtirish, ijodiy klublar, ekologik loyihalar",
            "5–7 sinf: loyiha va testlar, «Mening kelajakdagi kasbim» tayyorgarligi",
            "8–9 sinf: kasb yo‘naltirish, amaliy tajriba, texnikumga yo‘naltirish (50%)",
            "10–11 sinf: TOP-300 universitetlar, MOCK imtihonlar, ongli kasb tanlash",
            "Foydalanuvchilar: maslahatchi · bosh mutaxassis · direktor · ota-ona",
            "Xulosa: PF-86 va 343-PQ talablari bajariladi — bitta raqamli platforma, barcha maktablar",
        ],
        top=1.45,
        size=14,
    )
    add_footer(s7, prs, footer, "7 / 7")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Saved: {OUT}")


if __name__ == "__main__":
    build()
